# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "项目详细讲解_PDF解析重点_无乱码版.pptx"
FONT = "Microsoft YaHei"

SLIDES = [
    ("基于 RAG 的角色扮演系统项目讲解", ["重点模块：pdf_ingest_service.py", "主题：复杂 PDF 解析、向量入库、混合检索与对话生成", "适用场景：项目答辩 / 技术讲解 / 系统演示"], "title"),
    ("项目整体定位", ["本项目是一个面向角色对话的 RAG 系统。", "用户选择角色并提问，系统从该角色知识库中检索相关内容，再调用大模型生成回答。", "核心目标：让大模型回答基于上传 PDF/文本资料，而不是仅依赖模型自身记忆。", "重点解决：长文档、复杂 PDF、表格、扫描页、图片图表的知识入库与检索。"], "content"),
    ("系统总体架构", ["前端：角色选择、聊天界面、知识文件上传、来源展示。", "后端：FastAPI 接口层 + Service 业务层 + Repository 数据访问层。", "数据层：MySQL 保存用户/角色/会话，Redis 保存短期记忆，Milvus 保存向量知识，Neo4j 保存图谱关系。", "AI 能力：Embedding、Rerank、LLM 生成、视觉 / OCR 理解。"], "content"),
    ("一次完整问答链路", ["用户输入问题，前端发起聊天请求。", "ChatService 验证用户、角色和会话状态。", "读取短期记忆，并可选进行 Query Rewrite。", "PDFIngestService 执行 Hybrid 检索，召回知识片段。", "LLMService 拼接角色设定、记忆、实时上下文和 RAG 上下文生成回复。", "保存对话、来源和记忆，并返回给前端。"], "content"),
    ("为什么 pdf_ingest_service.py 是核心", ["它同时负责“知识入库”和“知识检索”。", "入库侧：PDF → 文本 / OCR / 表格 / 图片描述 → chunk → embedding → Milvus。", "检索侧：用户问题 → BM25 + 向量检索 + Neo4j 图谱 → 融合 → Rerank。", "它决定了系统能不能从复杂 PDF 中提取出可被大模型利用的知识。"], "content"),
    ("PDF 入库主流程：ingest_file()", ["入口函数：ingest_file(character_id, pdf_path)。", "检查 PDF 文件是否存在。", "调用 _extract_text() 抽取完整可检索文本。", "调用 _chunk_text() 按窗口切分文本。", "每个 chunk 调用 _build_row() 生成入库记录。", "调用 _insert_into_milvus() 批量写入向量库。"], "content"),
    ("复杂 PDF 解析总览：_extract_text()", ["优先使用 PyMuPDF，因为它能处理文本层、表格、图片和页面渲染。", "如果 PyMuPDF 不可用，则退化为 pypdf 的普通文本抽取。", "逐页解析 PDF，并把每页内容合并成统一文本。", "每页内容由四类信息组成：文本层、OCR 文本、Markdown 表格、图片描述。"], "content"),
    ("文本层解析", ["对正常 PDF，直接使用 page.get_text(\"text\") 抽取可复制文本。", "优点：速度快、准确率高、对普通文档最稳定。", "适合：正文、标题、段落、普通说明文字。", "如果抽取结果很少，系统会判断该页可能是扫描页或图片页。"], "content"),
    ("OCR 扫描件解析", ["触发条件：当前页文本层长度少于 30 个字符。", "OCR 引擎：RapidOCR，基于 ONNX Runtime，轻量且可本地运行。", "处理流程：PDF 页面 → 2 倍缩放渲染 → 图片 → RapidOCR → 文本行拼接。", "设计理由：普通页不 OCR，避免降低速度；扫描页才 OCR，作为兜底。", "容错：未安装 OCR 或识别失败时，不中断整个入库流程。"], "content"),
    ("表格解析：保留二维结构", ["函数：_extract_tables_as_markdown(page)。", "使用 PyMuPDF 的 page.find_tables() 自动检测表格区域。", "table.extract() 得到二维数组，再转换成 Markdown 表格。", "相比普通文本抽取，Markdown 能保留“表头—行—列”的对应关系。", "适合财务报表、股权结构表、客户销售表、募集资金用途表等。"], "content"),
    ("表格为什么对 RAG 很重要", ["普通 PDF 抽取可能打乱行列，导致“年份、指标、金额”关系丢失。", "Markdown 表格能让检索和大模型更容易理解数据归属。", "例如：2022 年营业收入、净利润、毛利率可以保持在同一行上下文中。", "对招股说明书、财报类问答尤其关键。"], "content"),
    ("图片与图表解析", ["函数：_extract_images_as_text(page)。", "系统会枚举 PDF 页面中的图片，并通过 xref 提取图片二进制。", "过滤规则：每页最多处理前 3 张，且宽高都不小于 100 像素。", "小图标、Logo、装饰元素会被跳过，避免浪费视觉模型调用成本。", "大图、组织结构图、趋势图、截图会送入视觉模型生成中文描述。"], "content"),
    ("视觉模型如何把图片变成知识", ["函数：_describe_image(image_bytes, ext)。", "图片会被 Base64 编码成 data URL，发送到 OpenAI 兼容 chat/completions 接口。", "使用配置项 vision_model_name，当前为 deepseek-ai/DeepSeek-OCR。", "提示词要求模型用中文描述图像中的关键信息，如图表类型、数据趋势、组织结构。", "最终格式：[图像内容描述] + 视觉模型输出，并追加到页面文本中。"], "content"),
    ("文本切分策略：_chunk_text()", ["默认 chunk_size = 800 字符，overlap = 120 字符。", "800 字符能兼顾段落完整性和检索粒度。", "120 字符重叠可以避免股东名称、金额、比例等信息被切在边界。", "切分前会压缩连续空白，减少无意义字符对 embedding 的影响。", "结果：长 PDF 被拆成多个可检索、可向量化的文本块。"], "content"),
    ("每个 chunk 的入库结构", ["函数：_build_row()。", "source_file：来源 PDF 文件名，用于前端展示引用。", "chunk_index：文本块序号，用于定位原文。", "text：真正进入 RAG 上下文的文本。", "keywords：jieba 提取的关键词，用于 BM25 检索增强。", "vector：Embedding 向量，用于 Milvus 语义检索。", "chunk_hash：SHA256 指纹，用于去重和追踪。"], "content"),
    ("Embedding 与 Milvus 入库", ["Embedding API：调用 OpenAI 兼容 /embeddings 接口。", "模型名来自 settings.embedding_model_name，向量维度来自 settings.milvus_dim。", "当前配置维度为 1024，适配 bge-large-zh-v1.5 / bge-m3 等 1024 维模型。", "Embedding 有 LRU 缓存，最多缓存 512 条，减少重复调用。", "API 失败时回退到 SHA256 伪向量，保证流程不中断。"], "content"),
    ("Milvus Collection 设计", ["每个角色一个独立 collection：character_knowledge_{id}。", "隔离不同角色知识，避免跨角色知识污染。", "字段包括：id、source_file、chunk_index、chunk_hash、text、keywords、vector。", "vector 字段使用 FLOAT_VECTOR，维度由 settings.milvus_dim 控制。", "索引使用 IVF_FLAT，距离度量使用 COSINE。"], "content"),
    ("BM25 关键词检索", ["函数：search_keyword()。", "作用：补充向量检索在精确数字、年份、金额、人名、公司名上的不足。", "数据来源：从 Milvus 拉取 text + keywords，并在 Python 内存中计算 BM25。", "参数：k1 = 1.2，b = 0.75。", "适合问题：2022 年营业收入是多少？控股股东是谁？发行数量是多少？"], "content"),
    ("向量检索与语义召回", ["函数：search_vector()。", "先把用户问题也转成 embedding 向量。", "使用 Milvus 在 vector 字段上做 ANN 检索。", "度量方式为 COSINE，适合语义相似度。", "优势：能处理同义改写和自然语言表达，例如“主营业务”与“主要从事”。"], "content"),
    ("Hybrid 检索融合", ["函数：search_hybrid()。", "三路召回：BM25 关键词检索、Milvus 向量检索、Neo4j 图谱检索。", "BM25 强在精确匹配，向量强在语义泛化，Neo4j 强在实体关系。", "融合权重：向量 0.6，关键词 0.4；图谱作为结构化候选补充。", "合并后按 hybrid_score 排序，并扩大候选池给 Rerank。"], "content"),
    ("Rerank 精排", ["函数：_rerank()。", "第一阶段 Hybrid 检索强调“召回”，宁可多取候选。", "Rerank 模型逐对判断“问题—片段”的真实相关性。", "当前模型配置：BAAI/bge-reranker-v2-m3。", "Rerank 失败时回退到 hybrid_score 排序，保证系统可用性。"], "content"),
    ("容错与稳定性设计", ["PyMuPDF 不可用 → 回退 pypdf。", "OCR 未安装或失败 → 跳过 OCR，不影响普通文本。", "表格检测失败 → 跳过该表格。", "图片描述失败 → 跳过图片描述。", "Embedding API 失败 → 回退 SHA256 伪向量。", "Neo4j / Rerank 失败 → 返回空图谱或使用融合分回退。"], "content"),
    ("与 ChatService 的连接", ["ChatService 是问答编排中心。", "接收问题后读取记忆、进行 Query Rewrite，再调用 PDFIngestService 检索上下文。", "检索结果 context 会和角色设定、记忆、实时上下文一起送入 LLM。", "最终答案、RAG 来源、对话记录会写回数据库。", "流式接口中检索时会先输出“正在进行检索，请稍等……”降低用户感知延迟。"], "content"),
    ("项目亮点总结", ["多模态 PDF 入库：文本、OCR、表格、图片描述统一转成可检索文本。", "角色级知识隔离：每个角色拥有独立 Milvus collection。", "Hybrid 检索：BM25 + 向量 + Neo4j 兼顾精确事实、语义泛化和关系推理。", "Rerank 精排：提升最终上下文质量，降低幻觉风险。", "强容错设计：外部服务失败时自动降级，保证主流程可运行。"], "content"),
    ("可改进方向", ["引入更专业的表格解析工具，提升复杂跨页表格还原能力。", "对 OCR 结果增加置信度过滤和版面顺序重排。", "将 BM25 从内存计算升级到 Elasticsearch / OpenSearch，支持更大规模文档。", "补充异步入库队列，避免上传大 PDF 时阻塞请求。", "为 chunk 增加页码范围和表格/图片类型标签，方便前端引用定位。"], "content"),
    ("结束页", ["本项目的关键不是简单调用大模型，而是构建了一条完整的知识处理链路。", "pdf_ingest_service.py 将复杂 PDF 转化为结构化、可检索、可追溯的知识资产。", "最终让角色对话具备基于文档的事实回答能力。"], "end"),
]

BLUE = RGBColor(30, 58, 138)
DARK = RGBColor(15, 23, 42)
TEXT = RGBColor(51, 65, 85)
LIGHT = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(236, 253, 245)


def apply_font(paragraph, size=18, color=TEXT, bold=False):
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    for run in paragraph.runs:
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold


def add_text(slide, text, x, y, w, h, size=24, color=TEXT, bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    apply_font(p, size, color, bold)
    return box


def add_bullets(slide, items, x, y, w, h, size=18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(9)
        apply_font(p, size, TEXT, False)
    return box


def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for title, items, kind in SLIDES:
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = LIGHT

        if kind == "title":
            bg.fore_color.rgb = DARK
            add_text(slide, title, 0.8, 1.45, 11.8, 0.8, 34, WHITE, True)
            add_text(slide, items[0], 0.86, 2.55, 11.6, 0.5, 20, RGBColor(191, 219, 254), False)
            add_bullets(slide, items[1:], 1.05, 3.35, 11.0, 1.6, 18)
        elif kind == "end":
            add_text(slide, title, 0.8, 0.85, 11.8, 0.7, 32, DARK, True)
            add_rect(slide, 0.9, 1.95, 11.5, 3.3, GREEN, RGBColor(134, 239, 172))
            add_bullets(slide, items, 1.25, 2.35, 10.8, 2.45, 20)
        else:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.58))
            bar.fill.solid()
            bar.fill.fore_color.rgb = BLUE
            bar.line.color.rgb = BLUE
            add_text(slide, title, 0.55, 0.12, 12.2, 0.45, 24, WHITE, True)
            add_rect(slide, 0.7, 0.95, 11.95, 5.75, WHITE, RGBColor(226, 232, 240))
            add_bullets(slide, items, 1.08, 1.28, 11.1, 4.95, 18)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
