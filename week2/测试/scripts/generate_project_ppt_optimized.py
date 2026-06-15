# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "项目详细讲解_PDF解析重点_优化版.pptx"
FONT = "Microsoft YaHei"

NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(6, 182, 212)
GREEN = RGBColor(16, 185, 129)
AMBER = RGBColor(245, 158, 11)
PURPLE = RGBColor(124, 58, 237)
RED = RGBColor(239, 68, 68)
BG = RGBColor(248, 250, 252)
CARD = RGBColor(255, 255, 255)
TEXT = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
LINE = RGBColor(226, 232, 240)
WHITE = RGBColor(255, 255, 255)

SLIDES = [
    {"kind": "cover", "title": "基于 RAG 的角色扮演系统", "subtitle": "项目详细讲解 · PDF 解析与稳定性优化版", "tags": ["复杂 PDF 解析", "Hybrid 检索", "LLMClient 封装", "生产安全"]},
    {"kind": "section", "title": "01 项目整体理解", "subtitle": "先看系统做什么，再深入 PDF 入库与检索核心。"},
    {"kind": "cards", "title": "项目定位", "lead": "让角色对话具备基于文档的事实回答能力。", "cards": [("用户侧", "选择角色、上传知识、发起对话、查看引用来源", BLUE), ("知识侧", "PDF / 文本资料被解析、切分、向量化并写入 Milvus", GREEN), ("模型侧", "结合角色设定、记忆、实时上下文和 RAG 上下文生成回答", PURPLE)]},
    {"kind": "flow", "title": "一次完整问答链路", "steps": [("用户提问", BLUE), ("ChatService\n编排", CYAN), ("RAG 检索", GREEN), ("LLM 生成", PURPLE), ("保存对话\n返回来源", AMBER)], "note": "核心思想：先检索可靠上下文，再让大模型基于上下文回答。"},
    {"kind": "architecture", "title": "系统总体架构", "layers": [("前端交互层", ["角色选择", "聊天界面", "知识上传", "来源展示"], BLUE), ("后端业务层", ["FastAPI", "ChatService", "KnowledgeService", "PDFIngestService"], GREEN), ("数据与 AI 层", ["MySQL", "Redis", "Milvus", "Neo4j", "LLMClient"], PURPLE)]},
    {"kind": "section", "title": "02 PDFIngestService 核心链路", "subtitle": "它是项目中最关键的知识处理管道。"},
    {"kind": "split", "title": "为什么它是核心模块？", "left_title": "知识入库", "left": ["PDF 解析", "OCR 兜底", "表格转 Markdown", "图片转描述", "切分 + 向量化 + 入库"], "right_title": "知识检索", "right": ["BM25 关键词召回", "Milvus 向量召回", "Neo4j 图谱召回", "Hybrid 融合", "Rerank 精排"]},
    {"kind": "flow", "title": "PDF 入库主流程：ingest_file()", "steps": [("PDF 文件", BLUE), ("抽取文本\n_extract_text", CYAN), ("切分 Chunk\n_chunk_text", GREEN), ("生成记录\n_build_row", PURPLE), ("写入 Milvus\n_insert", AMBER)], "note": "一份 PDF 最终会变成多条带文本、关键词、向量和来源信息的知识记录。"},
    {"kind": "matrix", "title": "复杂 PDF 解析：四路信息合并", "items": [("文本层", "page.get_text 直接提取可复制文本", BLUE), ("OCR", "扫描页或图片页通过 RapidOCR 识别", GREEN), ("表格", "find_tables 后转 Markdown 保留行列结构", AMBER), ("图片", "视觉模型生成中文描述进入检索", PURPLE)]},
    {"kind": "highlight", "title": "OCR 扫描件解析", "headline": "只在需要时启用 OCR，兼顾速度与覆盖率", "points": ["触发条件：页面文本层少于 30 个字符。", "流程：PDF 页面 → 2 倍渲染 → 图片 → RapidOCR → 文本拼接。", "价值：扫描版 PDF、截图页、图片型资料也能进入知识库。", "容错：OCR 不可用时跳过，不影响普通 PDF 入库。"], "color": GREEN},
    {"kind": "highlight", "title": "表格解析：把二维结构保留下来", "headline": "财报、股权、募集资金等表格不能只当普通文本处理", "points": ["使用 PyMuPDF 的 page.find_tables() 检测表格。", "table.extract() 得到二维数组。", "转换成 Markdown 表格，保留表头、行、列关系。", "提升金额、年份、比例、指标类问题的回答稳定性。"], "color": AMBER},
    {"kind": "highlight", "title": "图片与图表解析", "headline": "让图表、组织结构图、截图也能被 RAG 检索", "points": ["每页最多处理前 3 张大图，过滤 Logo 和装饰元素。", "通过 xref 提取图片二进制。", "Base64 data URL 发送给视觉模型。", "生成中文图片描述并追加到页面文本中。"], "color": PURPLE},
    {"kind": "cards", "title": "Chunk 入库字段设计", "lead": "每个文本块不只是 text，而是一条可追溯、可检索、可融合的知识记录。", "cards": [("定位", "source_file + chunk_index：展示引用来源并定位原文", BLUE), ("检索", "keywords + vector：同时支持 BM25 与向量召回", GREEN), ("稳定", "chunk_hash：生成文本指纹，便于去重和追踪", AMBER)]},
    {"kind": "section", "title": "03 检索与回答生成", "subtitle": "从“能入库”到“能准确召回”。"},
    {"kind": "compare", "title": "为什么要 Hybrid 检索？", "cols": [("BM25", "擅长数字、年份、人名、公司名等精确匹配", AMBER), ("向量检索", "擅长语义泛化和同义改写，例如“主营业务”与“主要从事”", BLUE), ("Neo4j", "擅长股权关系、关联关系、实体之间的结构化连接", GREEN)]},
    {"kind": "flow", "title": "Hybrid 检索融合流程", "steps": [("问题", BLUE), ("BM25", AMBER), ("Vector", CYAN), ("Neo4j", GREEN), ("Rerank\nTop-K", PURPLE)], "note": "默认融合权重：向量 0.6，关键词 0.4；图谱作为结构化候选补充。"},
    {"kind": "highlight", "title": "Hybrid 检索稳定性优化", "headline": "三路召回并行执行，单路失败自动降级", "points": ["ThreadPoolExecutor 同时提交 BM25、Milvus Vector、Neo4j 三路召回。", "_safe_retrieval_result 统一接收 Future 结果，异常时记录日志并返回空列表。", "Milvus、BM25 或 Neo4j 任一路临时失败，不会中断整次问答。", "融合阶段继续使用其他可用候选，保证 RAG 主流程优先可用。"], "color": GREEN},
    {"kind": "highlight", "title": "Rerank 精排", "headline": "第一阶段负责召回，第二阶段负责排序", "points": ["Hybrid 先尽量多召回候选，避免漏掉答案。", "Rerank 判断“问题—片段”的真实相关性。", "减少无关上下文进入 LLM，降低幻觉概率。", "外部 Rerank 失败时回退到 hybrid_score，保证可用性。"], "color": PURPLE},
    {"kind": "architecture", "title": "容错与降级设计", "layers": [("解析降级", ["PyMuPDF 不可用 → pypdf", "OCR 失败 → 跳过", "表格失败 → 跳过单表"], GREEN), ("AI 降级", ["图片描述失败 → 跳过", "Embedding 失败 → SHA256 伪向量", "Rerank 失败 → 融合分排序"], AMBER), ("检索降级", ["单路异常隔离", "Neo4j 失败 → 空图谱", "缓存减少重复计算"], BLUE)]},
    {"kind": "section", "title": "04 LLM 调用与工程优化", "subtitle": "把外部模型调用、跨域和生产配置做成统一、可维护的工程能力。"},
    {"kind": "split", "title": "LLM 调用统一封装", "left_title": "LLMService", "left": ["负责 Prompt 编排", "支持 generate / generate_stream", "负责摘要 summarize", "负责 Query Rewrite", "解析 SSE 增量内容"], "right_title": "LLMClient", "right": ["统一拼接 /chat/completions", "统一生成 Authorization", "统一构造 chat_payload", "封装非流式 chat", "封装流式 chat_stream"]},
    {"kind": "highlight", "title": "为什么拆出 LLMClient？", "headline": "把 HTTP 细节从业务服务中剥离，避免重复和分散维护", "points": ["摘要、Query Rewrite、普通生成、流式生成都复用同一套 HTTP 调用逻辑。", "以后切换 OpenAI、SiliconFlow、vLLM、SGLang，只需调整 base_url、model、key。", "LLMService 专注业务语义：角色人设、RAG 上下文、记忆和实时上下文。", "LLMClient 专注传输细节：URL、Header、Payload、Timeout、响应解析。"], "color": BLUE},
    {"kind": "matrix", "title": "配置与安全优化", "items": [("CORS 白名单", "cors_allowed_origins 与 cors_allow_origin_regex 限制前端来源，避免生产环境无限放开跨域", BLUE), ("生产校验", "validate_runtime 检查 debug、JWT 密钥、Neo4j 默认密码、LLM API Key 等风险项", RED), ("集中配置", "Settings 统一管理 MySQL、Redis、Milvus、Neo4j、LLM、RAG 参数", GREEN), ("参数化检索", "retrieval_top_k、rerank_top_k、hybrid 权重均可通过配置调整", PURPLE)]},
    {"kind": "section", "title": "04 项目价值与优化方向", "subtitle": "总结项目亮点，并给出后续改进空间。"},
    {"kind": "cards", "title": "项目亮点总结", "lead": "这个项目的价值在于构建了完整的知识处理闭环。", "cards": [("多模态 PDF 入库", "文本、OCR、表格、图片描述统一进入知识库", GREEN), ("混合检索增强", "BM25 + Vector + Graph 并行召回，单路异常自动降级", BLUE), ("工程可维护性", "LLMClient 统一封装调用，配置校验提升上线安全", PURPLE)]},
    {"kind": "matrix", "title": "后续优化方向", "items": [("表格增强", "处理跨页表格、复杂合并单元格", AMBER), ("异步入库", "上传大 PDF 后后台队列处理", BLUE), ("统一客户端扩展", "Embedding / Rerank 后续也可抽象成独立 Client", GREEN), ("检索扩展", "BM25 可升级 Elasticsearch / OpenSearch", PURPLE)]},
    {"kind": "end", "title": "总结", "points": ["本项目不是简单调用大模型，而是构建了完整的 RAG 知识链路。", "pdf_ingest_service.py 负责把复杂 PDF 转成可检索、可追溯、可融合的知识资产。", "Hybrid 检索通过单路异常隔离提升可用性，LLMClient 通过统一封装提升可维护性。", "最终目标：让角色对话能够基于真实文档进行稳定、可信的事实回答。"]},
]


def add_transition(slide):
    transition = OxmlElement("p:transition")
    transition.set("spd", "med")
    transition.append(OxmlElement("p:fade"))
    slide._element.insert(-1, transition)


def set_font(paragraph, size, color=TEXT, bold=False):
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    for run in paragraph.runs:
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold


def add_text(slide, text, x, y, w, h, size=20, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    set_font(p, size, color, bold)
    return box


def add_body(slide, items, x, y, w, h, size=16):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        set_font(p, size, TEXT)
    return box


def rect(slide, x, y, w, h, fill, line=None, radius=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def pill(slide, text, x, y, w, color):
    shape = rect(slide, x, y, w, 0.36, color, None)
    add_text(slide, text, x, y + 0.02, w, 0.3, 11, WHITE, True, PP_ALIGN.CENTER)
    return shape


def header(slide, title, idx, total):
    rect(slide, 0, 0, 13.333, 0.66, NAVY, NAVY, False)
    add_text(slide, title, 0.55, 0.12, 9.7, 0.4, 21, WHITE, True)
    add_text(slide, f"{idx:02d} / {total:02d}", 11.5, 0.14, 1.25, 0.35, 11, RGBColor(203, 213, 225), False, PP_ALIGN.RIGHT)


def footer(slide):
    add_text(slide, "Role-playing system based on RAG · PDF Ingest Deep Dive", 0.55, 7.05, 7.0, 0.2, 8, MUTED)


def line(slide, x1, y1, x2, y2, color=LINE):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)
    return connector


def draw_cover(slide, data):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    rect(slide, 8.7, -0.35, 5.1, 8.2, RGBColor(30, 64, 175), None, False)
    rect(slide, 9.5, 0.8, 2.8, 2.8, RGBColor(59, 130, 246), None)
    rect(slide, 8.4, 4.7, 3.8, 1.7, RGBColor(14, 165, 233), None)
    add_text(slide, data["title"], 0.8, 1.55, 7.6, 0.85, 34, WHITE, True)
    add_text(slide, data["subtitle"], 0.85, 2.55, 7.6, 0.45, 19, RGBColor(191, 219, 254))
    x = 0.85
    for tag in data["tags"]:
        pill(slide, tag, x, 3.35, 1.8 if len(tag) < 8 else 2.25, BLUE)
        x += 2.05 if len(tag) < 8 else 2.5
    add_text(slide, "从文档解析到知识检索，再到可信问答生成", 0.85, 5.65, 7.0, 0.45, 16, RGBColor(226, 232, 240))


def draw_section(slide, data):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    rect(slide, 0.8, 1.45, 0.12, 3.2, BLUE, None, False)
    add_text(slide, data["title"], 1.15, 2.1, 10.5, 0.75, 34, WHITE, True)
    add_text(slide, data["subtitle"], 1.18, 3.0, 9.5, 0.5, 18, RGBColor(203, 213, 225))


def draw_cards(slide, data, idx, total):
    header(slide, data["title"], idx, total)
    footer(slide)
    add_text(slide, data["lead"], 0.75, 0.95, 11.9, 0.55, 20, NAVY, True)
    for i, (title, desc, color) in enumerate(data["cards"]):
        x = 0.85 + i * 4.15
        rect(slide, x, 1.85, 3.65, 3.95, CARD, LINE)
        rect(slide, x + 0.28, 2.18, 0.55, 0.55, color, None)
        add_text(slide, title, x + 0.28, 2.95, 3.1, 0.45, 20, NAVY, True)
        add_text(slide, desc, x + 0.28, 3.55, 3.05, 1.45, 15, TEXT)


def draw_flow(slide, data, idx, total):
    header(slide, data["title"], idx, total)
    footer(slide)
    y = 2.35
    w = 2.0
    gap = 0.42
    start = 0.72
    for i, (label, color) in enumerate(data["steps"]):
        x = start + i * (w + gap)
        rect(slide, x, y, w, 1.18, color, None)
        add_text(slide, label, x + 0.12, y + 0.2, w - 0.24, 0.72, 15, WHITE, True, PP_ALIGN.CENTER)
        if i < len(data["steps"]) - 1:
            line(slide, x + w, y + 0.58, x + w + gap, y + 0.58, RGBColor(148, 163, 184))
    rect(slide, 1.25, 4.55, 10.8, 1.05, RGBColor(239, 246, 255), RGBColor(191, 219, 254))
    add_text(slide, data["note"], 1.55, 4.78, 10.2, 0.45, 17, BLUE, True, PP_ALIGN.CENTER)


def draw_architecture(slide, data, idx, total):
    header(slide, data["title"], idx, total)
    footer(slide)
    for i, (title, items, color) in enumerate(data["layers"]):
        y = 1.12 + i * 1.75
        rect(slide, 0.85, y, 11.65, 1.25, CARD, LINE)
        rect(slide, 0.85, y, 2.45, 1.25, color, None)
        add_text(slide, title, 1.05, y + 0.34, 2.05, 0.42, 17, WHITE, True, PP_ALIGN.CENTER)
        x = 3.55
        for item in items:
            pill(slide, item, x, y + 0.43, 1.55 if len(item) <= 8 else 2.05, RGBColor(241, 245, 249))
            add_text(slide, item, x, y + 0.45, 1.55 if len(item) <= 8 else 2.05, 0.26, 10, TEXT, True, PP_ALIGN.CENTER)
            x += 1.75 if len(item) <= 8 else 2.25


def draw_split(slide, data, idx, total):
    header(slide, data["title"], idx, total)
    footer(slide)
    rect(slide, 0.9, 1.25, 5.55, 4.95, RGBColor(239, 246, 255), RGBColor(191, 219, 254))
    rect(slide, 6.9, 1.25, 5.55, 4.95, RGBColor(250, 245, 255), RGBColor(216, 180, 254))
    add_text(slide, data["left_title"], 1.25, 1.62, 4.7, 0.42, 23, BLUE, True, PP_ALIGN.CENTER)
    add_text(slide, data["right_title"], 7.25, 1.62, 4.7, 0.42, 23, PURPLE, True, PP_ALIGN.CENTER)
    add_body(slide, data["left"], 1.35, 2.35, 4.75, 2.9, 17)
    add_body(slide, data["right"], 7.35, 2.35, 4.75, 2.9, 17)


def draw_matrix(slide, data, idx, total):
    header(slide, data["title"], idx, total)
    footer(slide)
    for i, (title, desc, color) in enumerate(data["items"]):
        x = 0.95 + (i % 2) * 6.05
        y = 1.25 + (i // 2) * 2.45
        rect(slide, x, y, 5.55, 1.9, CARD, LINE)
        rect(slide, x, y, 0.18, 1.9, color, None, False)
        add_text(slide, title, x + 0.45, y + 0.3, 4.7, 0.35, 21, NAVY, True)
        add_text(slide, desc, x + 0.45, y + 0.88, 4.65, 0.62, 15, TEXT)


def draw_highlight(slide, data, idx, total):
    header(slide, data["title"], idx, total)
    footer(slide)
    rect(slide, 0.9, 1.1, 11.55, 1.25, data["color"], None)
    add_text(slide, data["headline"], 1.2, 1.43, 10.9, 0.42, 22, WHITE, True, PP_ALIGN.CENTER)
    for i, point in enumerate(data["points"]):
        y = 2.82 + i * 0.82
        rect(slide, 1.25, y, 10.85, 0.58, CARD, LINE)
        rect(slide, 1.48, y + 0.15, 0.28, 0.28, data["color"], None)
        add_text(slide, point, 1.95, y + 0.08, 9.6, 0.32, 15, TEXT)


def draw_compare(slide, data, idx, total):
    header(slide, data["title"], idx, total)
    footer(slide)
    for i, (title, desc, color) in enumerate(data["cols"]):
        x = 0.85 + i * 4.15
        rect(slide, x, 1.35, 3.65, 4.85, CARD, LINE)
        rect(slide, x, 1.35, 3.65, 0.78, color, None)
        add_text(slide, title, x + 0.25, 1.55, 3.15, 0.32, 21, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, desc, x + 0.35, 2.55, 2.95, 2.2, 17, TEXT, False, PP_ALIGN.CENTER)


def draw_end(slide, data, idx, total):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_text(slide, data["title"], 0.85, 0.95, 11.5, 0.8, 34, WHITE, True, PP_ALIGN.CENTER)
    rect(slide, 1.25, 2.05, 10.85, 3.3, RGBColor(30, 41, 59), RGBColor(51, 65, 85))
    add_body(slide, data["points"], 1.85, 2.55, 9.75, 2.1, 19)
    add_text(slide, f"{idx:02d} / {total:02d}", 11.25, 6.95, 1.2, 0.25, 10, RGBColor(203, 213, 225), False, PP_ALIGN.RIGHT)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = len(SLIDES)

    for idx, data in enumerate(SLIDES, 1):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG
        add_transition(slide)
        kind = data["kind"]
        if kind == "cover":
            draw_cover(slide, data)
        elif kind == "section":
            draw_section(slide, data)
        elif kind == "cards":
            draw_cards(slide, data, idx, total)
        elif kind == "flow":
            draw_flow(slide, data, idx, total)
        elif kind == "architecture":
            draw_architecture(slide, data, idx, total)
        elif kind == "split":
            draw_split(slide, data, idx, total)
        elif kind == "matrix":
            draw_matrix(slide, data, idx, total)
        elif kind == "highlight":
            draw_highlight(slide, data, idx, total)
        elif kind == "compare":
            draw_compare(slide, data, idx, total)
        elif kind == "end":
            draw_end(slide, data, idx, total)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
