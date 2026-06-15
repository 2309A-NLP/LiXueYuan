from __future__ import annotations

import shutil
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
PPT = ROOT / "docs" / "工单" / "答辩.pptx"
BACKUP = ROOT / "docs" / "工单" / "答辩_原始备份.pptx"

TITLE = RGBColor(24, 45, 92)
ACCENT = RGBColor(37, 99, 235)
DARK = RGBColor(31, 41, 55)
MUTED = RGBColor(107, 114, 128)
LIGHT = RGBColor(239, 246, 255)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(16, 185, 129)
ORANGE = RGBColor(245, 158, 11)
PURPLE = RGBColor(124, 58, 237)


def set_text_frame(tf, paragraphs, font_size=18, color=DARK, bold_first=False):
    tf.clear()
    for idx, text in enumerate(paragraphs):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        if bold_first and idx == 0:
            p.font.bold = True


def add_textbox(slide, x, y, w, h, text, size=18, color=DARK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    if align:
        p.alignment = align
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.55, 0.25, 12.2, 0.45, title, 24, TITLE, True)
    if subtitle:
        add_textbox(slide, 0.58, 0.74, 12, 0.3, subtitle, 10.5, MUTED)
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.08), Inches(12.2), Inches(0.025))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.color.rgb = ACCENT


def add_card(slide, x, y, w, h, title, body, color=LIGHT):
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(219, 234, 254)
    add_textbox(slide, x + 0.18, y + 0.12, w - 0.36, 0.28, title, 14, TITLE, True)
    add_textbox(slide, x + 0.18, y + 0.48, w - 0.36, h - 0.58, body, 11.5, DARK)
    return shape


def add_bullets(slide, x, y, w, h, title, bullets, accent=ACCENT):
    add_textbox(slide, x, y, w, 0.35, title, 17, TITLE, True)
    box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.45), Inches(w), Inches(h - 0.45))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)
    return box


def add_footer(slide, idx):
    add_textbox(slide, 10.9, 7.08, 1.8, 0.22, f"{idx:02d}", 9, MUTED, False, PP_ALIGN.RIGHT)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


_PPTX_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def add_transition(slide, style: str = "fade", dur: int = 600):
    """给幻灯片添加切换特效。
    style: fade | push_left | push_up | wipe | split | zoom
    dur: 持续时间（毫秒）
    """
    elem = slide._element
    existing = elem.find(qn("p:transition"))
    if existing is not None:
        elem.remove(existing)

    speed_attr = f' dur="{dur}"'
    inner = {
        "fade":      "<p:fade/>",
        "push_left": '<p:push dir="l"/>',
        "push_up":   '<p:push dir="u"/>',
        "wipe":      '<p:wipe dir="l"/>',
        "split":     '<p:split dir="horz" orient="out"/>',
        "zoom":      "<p:zoom/>",
        "reveal":    '<p:reveal dir="l"/>',
        "cover":     '<p:cover dir="l"/>',
    }.get(style, "<p:fade/>")

    xml = (
        f'<p:transition xmlns:p="{_PPTX_NS}"{speed_attr} spd="med">'
        f'{inner}'
        f"</p:transition>"
    )
    elem.append(etree.fromstring(xml))


def build_ppt():
    if PPT.exists() and not BACKUP.exists():
        shutil.copy2(PPT, BACKUP)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ------------------------------------------------------------------ #
    # SLIDE 1  封面                                                         #
    # ------------------------------------------------------------------ #
    s = blank_slide(prs)
    bg = s.background.fill
    bg.solid(); bg.fore_color.rgb = RGBColor(248, 250, 252)
    add_textbox(s, 0.75, 1.05, 12, 0.65, "基于 RAG 的角色扮演问答系统", 34, TITLE, True)
    add_textbox(s, 0.8, 1.85, 10.8, 0.4,
                "让大模型真正读懂你上传的文件，基于文件内容回答问题，不胡编、有出处", 17, ACCENT, True)
    add_card(s, 0.85, 2.65, 3.8, 1.35, "核心目标",
             "上传一份 PDF 招股说明书，系统就能基于其中的真实内容回答财务、股权、风险等各类问题。")
    add_card(s, 4.9, 2.65, 3.8, 1.35, "整体流程",
             "上传文档 → 自动切段、存入检索库 → 用户提问 → 搜索相关段落 → 交给大模型生成答案")
    add_card(s, 8.95, 2.65, 3.55, 1.35, "典型用途",
             "招股说明书问答、财务数据查询、股权关系分析、角色扮演知识问答。")
    add_textbox(s, 0.85, 6.55, 5.8, 0.35, "汇报人：2309A_nlp-李学远", 14, MUTED)
    add_footer(s, 1)
    add_transition(s, "fade", 700)

    # ------------------------------------------------------------------ #
    # SLIDE 2  目录                                                         #
    # ------------------------------------------------------------------ #
    s = blank_slide(prs)
    add_title(s, "目录", "按照「为什么做、怎么做、做了什么、效果如何」四个层次展开")
    items = [
        ("01", "研究背景与问题",   "直接问大模型会怎样？为什么要引入检索增强？"),
        ("02", "系统架构与主流程", "各部分怎么分工，一条完整的回答是怎么产生的？"),
        ("03", "核心技术实现",     "文件怎么解析？问题怎么搜索？答案怎么生成？"),
        ("04", "工单演进与创新点", "系统怎么一步步从基础版演进到现在的能力？"),
        ("05", "测试评估与总结",   "系统效果怎么样？还有哪些可以继续改进的地方？"),
    ]
    y = 1.35
    for num, t, b in items:
        add_textbox(s, 0.9,  y,        0.65, 0.35, num, 18, ACCENT, True)
        add_textbox(s, 1.65, y,        3.4,  0.35, t,   17, TITLE,  True)
        add_textbox(s, 5.25, y + 0.03, 7.3,  0.35, b,   13, DARK)
        y += 0.9
    add_footer(s, 2)
    add_transition(s, "push_left", 500)

    # ------------------------------------------------------------------ #
    # SLIDE 3  项目背景                                                     #
    # ------------------------------------------------------------------ #
    s = blank_slide(prs)
    add_title(s, "为什么不直接问 ChatGPT？", "普通大模型不知道你私有文件里的内容，只能靠自己的记忆猜测")
    add_card(s, 0.75, 1.35, 3.7, 1.35, "❌ 直接问大模型的问题",
             "大模型不知道你上传的文件；\n它可能会「编」一个看起来合理的答案（称为「幻觉」）；\n无法告诉你答案来自哪一页哪一段。",
             RGBColor(254, 242, 242))
    add_card(s, 4.85, 1.35, 3.7, 1.35, "✅ 本系统的做法（RAG）",
             "先在文件里搜索和问题最相关的段落；\n把搜到的段落作为「参考资料」交给大模型；\n大模型只基于参考资料回答，有据可查。")
    add_card(s, 8.95, 1.35, 3.65, 1.35, "🎯 本项目目标",
             "让角色扮演系统能够基于上传的 PDF 文件，对金融、股权、财务等专业问题进行可信、有来源的回答。",
             RGBColor(240, 253, 244))
    add_textbox(s, 0.9, 3.5, 11.8, 0.55,
                "用户提问  →  在文件里搜索相关段落  →  智能排序筛选  →  大模型基于段落回答  →  返回结果和出处",
                20, TITLE, True, PP_ALIGN.CENTER)
    add_bullets(s, 1.0, 4.7, 11.4, 1.45, "这样做有哪些好处",
                ["减少「胡编」：大模型必须基于检索到的内容回答，无法凭空捏造。",
                 "支持私有文件：上传招股说明书后立刻就能针对它提问。",
                 "答案可追溯：每次回答都可以告诉用户依据来自哪个段落。"])
    add_footer(s, 3)
    add_transition(s, "push_left", 500)

    # ------------------------------------------------------------------ #
    # SLIDE 4  系统架构                                                     #
    # ------------------------------------------------------------------ #
    s = blank_slide(prs)
    add_title(s, "系统总体架构", "像流水线一样，每一层只做自己的事，出了问题互不影响")
    add_card(s, 0.55, 1.35, 2.1, 1.05, "用户界面（前端）",
             "用户聊天、上传文件、选择角色的操作页面", RGBColor(239, 246, 255))
    add_card(s, 2.95, 1.35, 2.1, 1.05, "接口层（后端入口）",
             "接收请求、验证身份、把结果实时「打字」传回前端")
    add_card(s, 5.35, 1.35, 2.25, 1.05, "对话协调层",
             "把记忆、检索、大模型调用串联起来编排出完整回答")
    add_card(s, 7.9, 1.35, 2.25, 1.05, "检索增强管道",
             "负责解析文件、存入检索库、搜索相关段落、精排结果")
    add_card(s, 10.45, 1.35, 2.2, 1.05, "AI 模型服务",
             "大模型生成回答 / 图片文字识别 / 向量化 / 精排模型")
    add_textbox(s, 0.75, 3.05, 12, 0.3,
                "底层存储：MySQL（用户与对话记录）｜Redis（短期记忆缓存）｜Milvus（向量检索库）｜Neo4j（关系图谱）｜Docker（一键部署）",
                14, TITLE, True, PP_ALIGN.CENTER)
    add_bullets(s, 0.9, 4.0, 5.6, 2.1, "各层做什么",
                ["用户界面：用户看到的聊天页面和文件上传功能。",
                 "对话协调层：把「记住历史对话、搜索文件、生成回答」这三件事串起来。",
                 "检索增强管道：把 PDF 处理成可搜索的片段，再根据问题找到最相关的片段。",
                 "AI 模型服务：大模型负责生成回答，精排模型负责重新排序搜索结果。"])
    add_bullets(s, 7.0, 4.0, 5.3, 2.1, "为什么分层设计",
                ["每层只负责自己的职责，逻辑清晰，易于维护。",
                 "某个组件出错时，其他部分尽量继续正常工作。",
                 "想换大模型、换搜索引擎，只需改对应那一层，不影响整体。"])
    add_footer(s, 4)
    add_transition(s, "wipe", 500)

    # ------------------------------------------------------------------ #
    # SLIDE 5  知识入库流程                                                 #
    # ------------------------------------------------------------------ #
    s = blank_slide(prs)
    add_title(s, "核心流程一：把文件变成可检索的知识", "上传 PDF 后，系统自动把里面的各种内容处理成可以被搜索的片段")
    steps = [
        "① 上传 PDF",
        "② 提取文字\n表格 / 图片 / 扫描页",
        "③ 切成小段\n（含重叠防止断裂）",
        "④ 文字转「数字指纹」\n便于计算机比较相似度",
        "⑤ 存入向量库\nMilvus",
        "⑥ 抽取实体关系\n存入图谱库"
    ]
    x = 0.55
    for i, st in enumerate(steps):
        add_card(s, x, 1.45, 1.8, 1.05, f"{i + 1}", st, RGBColor(239, 246, 255))
        x += 2.1
    add_bullets(s, 0.85, 3.15, 5.8, 2.6, "如何处理复杂 PDF",
                ["普通正文：直接逐页提取文字。",
                 "财务表格：自动识别表格结构，转成行列整齐的文本，防止数据乱序。",
                 "扫描件（图片化页面）：用 OCR 技术识别图片中的文字。",
                 "图表 / 股权结构图：用视觉 AI 生成一段中文描述，再作为文字存入检索库。"])
    add_bullets(s, 7.0, 3.15, 5.3, 2.6, "为什么这样做",
                ["统一格式：不管是文字、表格还是图片，最终都变成文字片段，统一被搜索。",
                 "防止信息丢失：招股说明书里的财务表格和股权图非常关键，不能跳过。",
                 "切段带重叠：每段之间有部分内容重叠，防止一句话被切断后找不到完整语义。"])
    add_footer(s, 5)
    add_transition(s, "push_left", 500)

    # ------------------------------------------------------------------ #
    # SLIDE 6  用户问答流程                                                 #
    # ------------------------------------------------------------------ #
    s = blank_slide(prs)
    add_title(s, "核心流程二：用户提问到得到答案", "系统先搜索文件、再结合记忆、最后让大模型基于证据回答，不是直接猜测")
    flow = ("用户输入问题  →  读取本次对话的历史记录  →  理解「它」「该公司」等指代"
            "  →  三路并行搜索文件  →  智能重新排序结果  →  组合参考资料和角色人设"
            "  →  大模型逐字生成回答  →  保存本轮对话")
    add_textbox(s, 0.75, 1.35, 11.9, 0.75, flow, 16, TITLE, True, PP_ALIGN.CENTER)
    add_bullets(s, 0.85, 2.65, 5.75, 3.0, "多轮对话是怎么实现的",
                ["系统会记住最近几轮对话（存在内存缓存里，速度快）。",
                 "如果用户说「它的净利润呢」，系统会自动把「它」补全为上一轮提到的公司名。",
                 "对话结束后写入数据库，下次打开还能看到历史记录。"])
    add_bullets(s, 7.0, 2.65, 5.25, 3.0, "回答为什么可信",
                ["大模型看不到它「不知道」的内容，只能基于我们搜索到的段落回答。",
                 "每个角色有自己独立的知识库，不同角色的文件不会互相干扰。",
                 "回答来源可以追溯到具体段落，降低「胡编」的风险。"])
    add_footer(s, 6)
    add_transition(s, "push_left", 500)

    # ------------------------------------------------------------------ #
    # SLIDE 7  混合检索 + 精排                                              #
    # ------------------------------------------------------------------ #
    s = blank_slide(prs)
    add_title(s, "核心技术：三路并行搜索 + 智能重新排序", "一种搜索方式不够用，所以同时用三种方式搜索，最后统一打分排序")
    add_card(s, 0.75, 1.35, 3.55, 1.5, "① 关键词精确搜索",
             "像在书里 Ctrl+F 一样，适合搜索公司名、年份、金额、股权比例等精确信息。",
             RGBColor(255, 251, 235))
    add_card(s, 4.9, 1.35, 3.55, 1.5, "② 语义相似度搜索",
             "把文字转成数字向量，通过计算距离找「意思相近」的段落，适合同义词和自然语言问法。")
    add_card(s, 9.05, 1.35, 3.25, 1.5, "③ 关系图谱搜索",
             "在知识关系图里查找，适合「谁控股谁、关联方是谁」这类结构关系问题。",
             RGBColor(245, 243, 255))
    add_textbox(s, 1.0, 3.5, 11.3, 0.45,
                "三路结果汇总 → 统一打分 → 本地智能精排模型重新排序 → 选出最相关的片段交给大模型",
                20, ACCENT, True, PP_ALIGN.CENTER)
    add_bullets(s, 1.1, 4.5, 10.9, 1.7, "这样做有什么好处",
                ["单靠语义搜索找不准数字和人名，单靠关键词又不理解同义词——三路结合覆盖更全面。",
                 "智能精排模型（bge-reranker）再次判断每个候选段落和问题的真实相关性，比简单打分更准。",
                 "精排模型部署在本地，不依赖网络，速度稳定、答辩演示不受影响。"])
    add_footer(s, 7)
    add_transition(s, "split", 500)

    # ------------------------------------------------------------------ #
    # SLIDE 8  13 个工单演进                                                #
    # ------------------------------------------------------------------ #
    s = blank_slide(prs)
    add_title(s, "13 个工单演进路线", "先让系统能跑起来，再一步步解决真实场景中遇到的问题")
    tasks_info = [
        ("01 基础 PDF 问答",  "✅ 已完成", RGBColor(239, 246, 255)),
        ("02 检索优化",       "✅ 已完成", RGBColor(239, 246, 255)),
        ("03 表格解析",       "✅ 已完成", RGBColor(239, 246, 255)),
        ("04 图像解析",       "✅ 已完成", RGBColor(239, 246, 255)),
        ("05 问题理解优化",   "✅ 已完成", RGBColor(240, 253, 244)),
        ("06 混合检索",       "✅ 已完成", RGBColor(240, 253, 244)),
        ("07 功能评估",       "🧪 脚本评估", RGBColor(240, 253, 244)),
        ("08 图谱问答（基础）","✅ 已完成", RGBColor(240, 253, 244)),
        ("09 图谱优化",       "🟡 部分完成", RGBColor(245, 243, 255)),
        ("10 容器化部署",     "✅ 已完成", RGBColor(245, 243, 255)),
        ("11 向量模型微调",   "📌 方案设计", RGBColor(245, 243, 255)),
        ("12 轻量图谱 RAG",   "🧪 实验完成", RGBColor(245, 243, 255)),
        ("13 性能优化",       "✅ 部分完成", RGBColor(245, 243, 255)),
    ]
    for i, (title, status, color) in enumerate(tasks_info):
        row = i // 4
        col = i % 4
        cx = 0.65 + col * 3.15
        cy = 1.3 + row * 1.48
        add_card(s, cx, cy, 2.85, 1.0, title, status, color)
    add_textbox(s, 0.8, 6.05, 11.8, 0.45,
                "演进思路：先跑通基础问答 → 优化复杂文档解析 → 强化检索能力 → 引入知识图谱 → 容器化部署 → 性能调优",
                14, TITLE, True, PP_ALIGN.CENTER)
    add_footer(s, 8)
    add_transition(s, "push_left", 500)

    # ------------------------------------------------------------------ #
    # SLIDE 9  技术难点                                                     #
    # ------------------------------------------------------------------ #
    s = blank_slide(prs)
    add_title(s, "六大技术难点与解决方案", "每个难点都来自真实使用场景中遇到的问题")
    add_card(s, 0.7, 1.25, 3.8, 1.4, "难点 1：PDF 内容五花八门",
             "文字、表格、扫描图片、流程图混在一起，普通抽取会丢很多信息。\n→ 用多种工具分别处理：文字直接提取，表格转整齐文本，图片用 AI 识别。",
             RGBColor(239, 246, 255))
    add_card(s, 4.85, 1.25, 3.8, 1.4, "难点 2：数字和人名搜不准",
             "语义搜索擅长理解意思，但对具体数字、年份、比例往往不够敏感。\n→ 同时用关键词精确搜索来弥补，两路结合取长补短。",
             RGBColor(255, 251, 235))
    add_card(s, 9.0, 1.25, 3.55, 1.4, "难点 3：找不到谁控股谁",
             "股东、关联方、控制关系靠普通搜索很难覆盖。\n→ 把这些关系单独存在知识图谱里，专门用来回答「关系型」问题。",
             RGBColor(245, 243, 255))
    add_card(s, 0.7, 3.35, 3.8, 1.4, "难点 4：用户说「它」指谁",
             "多轮对话中用户常省略主语，比如「它的净利润呢？」，系统要能理解指代对象。\n→ 系统记住历史对话，自动把省略的内容补全后再去搜索。",
             RGBColor(240, 253, 244))
    add_card(s, 4.85, 3.35, 3.8, 1.4, "难点 5：搜出来的结果排序不好",
             "三路搜索汇总后，候选段落的相关性参差不齐，不能直接用。\n→ 用专门的精排 AI 模型重新给每个候选段落打分排序，选出真正相关的。",
             RGBColor(239, 246, 255))
    add_card(s, 9.0, 3.35, 3.55, 1.4, "难点 6：多个组件怎么统一部署",
             "系统用到数据库、缓存、向量库、图谱库、前后端，依赖复杂。\n→ 用 Docker 容器化技术一键启动所有服务，环境统一不出错。",
             RGBColor(240, 253, 244))
    add_footer(s, 9)
    add_transition(s, "push_left", 500)

    # ------------------------------------------------------------------ #
    # SLIDE 10  测试评估                                                    #
    # ------------------------------------------------------------------ #
    s = blank_slide(prs)
    add_title(s, "测试评估与性能优化", "用 15 道标准问题测试，从准确率、召回率、响应速度三个维度衡量效果")
    add_bullets(s, 0.8, 1.35, 5.8, 2.35, "怎么评估效果",
                ["搜索命中率：关键答案相关的段落有没有被搜索到？",
                 "段落相关性：搜索到的段落是否真的能回答这道题？",
                 "回答忠实性：大模型的回答是否基于搜索内容，有没有捏造？",
                 "响应速度：从提问到返回答案，每个环节各花了多少时间？"])
    add_bullets(s, 7.0, 1.35, 5.3, 2.35, "做了哪些速度优化",
                ["文档缓存：把常用段落提前缓存，避免每次提问都重新从库里拉数据。",
                 "三路搜索并行：三种搜索方式同时跑，不排队等待，节省时间。",
                 "先扩大候选再精排：先多搜一些候选，精排后再缩减，避免遗漏好答案。",
                 "精排模型本地化：精排不依赖外部网络，响应稳定。"])
    add_textbox(s, 0.95, 5.1, 11.4, 0.55,
                "结论：RAG 系统相比直接问大模型，在私有文档问答场景中准确率显著更高，且每个回答都有文档来源可查。",
                18, ACCENT, True, PP_ALIGN.CENTER)
    add_footer(s, 10)
    add_transition(s, "wipe", 500)

    # ------------------------------------------------------------------ #
    # SLIDE 11  系统部署                                                    #
    # ------------------------------------------------------------------ #
    s = blank_slide(prs)
    add_title(s, "系统部署与工程化", "一键启动所有服务，关机数据不丢，配置灵活可切换")
    add_card(s, 0.8, 1.35, 3.6, 1.3, "🐳 容器化一键部署",
             "用 Docker Compose 一条命令同时启动数据库、缓存、向量库、图谱库、前后端服务。")
    add_card(s, 4.9, 1.35, 3.6, 1.3, "💾 数据不会随关机丢失",
             "所有数据库都配置了持久化存储（磁盘挂载），关机重启后对话记录和知识库完整保留。",
             RGBColor(240, 253, 244))
    add_card(s, 9.0, 1.35, 3.3, 1.3, "⚙️ 灵活配置",
             "通过配置文件和环境变量管理所有参数，不用改代码就能切换模型、数据库地址和检索策略。",
             RGBColor(245, 243, 255))
    add_bullets(s, 1.0, 3.45, 11.0, 2.2, "工程设计亮点",
                ["每个角色的知识库独立隔离，不同角色之间的文件内容互不干扰。",
                 "某个组件（如图谱服务）出现故障时，系统尽量回退到其他检索方式，不直接崩溃中断回答。",
                 "前端回答边生成边显示（像打字一样），用户不用等到全部生成完才能看到内容。",
                 "大模型、向量化模型、精排模型都可以通过配置或少量代码切换，便于后续升级。"])
    add_footer(s, 11)
    add_transition(s, "push_left", 500)

    # ------------------------------------------------------------------ #
    # SLIDE 12  创新点                                                      #
    # ------------------------------------------------------------------ #
    s = blank_slide(prs)
    add_title(s, "项目创新点", "不只是简单调用大模型，而是围绕可信问答构建了完整的工程体系")
    add_card(s, 0.75, 1.25, 3.55, 1.3, "创新 1：多类型文档统一入库",
             "文字、表格、OCR 识别文字、图片 AI 描述——四类内容统一处理后进入同一个检索库。")
    add_card(s, 4.9, 1.25, 3.55, 1.3, "创新 2：三路融合搜索",
             "关键词精确搜索 + 语义相似度搜索 + 实体关系图谱搜索，三路结果合并打分，覆盖更全面。")
    add_card(s, 9.05, 1.25, 3.25, 1.3, "创新 3：本地精排模型",
             "搜索结果由本地 AI 精排模型重新打分排序，不依赖网络，速度稳定，效果更好。")
    add_card(s, 0.75, 3.2, 3.55, 1.3, "创新 4：角色有记忆",
             "系统同时记住短期对话（速度快）和长期历史（持久存档），支持自然的多轮连续对话。")
    add_card(s, 4.9, 3.2, 3.55, 1.3, "创新 5：知识图谱辅助",
             "用知识关系图补充「谁控股谁、关联方是谁」等结构化关系问题，纯文字搜索覆盖不到这类问题。")
    add_card(s, 9.05, 3.2, 3.25, 1.3, "创新 6：完整工程闭环",
             "从文件上传、检索、回答生成、效果评估到容器化部署，形成完整可运行的系统。")
    add_footer(s, 12)
    add_transition(s, "split", 500)

    # ------------------------------------------------------------------ #
    # SLIDE 13  总结与展望                                                  #
    # ------------------------------------------------------------------ #
    s = blank_slide(prs)
    add_title(s, "总结与展望", "系统已能实现完整的文档知识问答，后续仍有值得继续深入的方向")
    add_bullets(s, 0.85, 1.35, 5.8, 2.5, "目前已经实现",
                ["上传 PDF，自动解析文字、表格、图片并存入检索库。",
                 "用户提问后，系统搜索相关段落，让大模型基于段落生成回答。",
                 "系统记住多轮对话，能理解上下文中「它」「该公司」等指代表达。",
                 "三路并行搜索 + 本地精排模型，提升搜索准确率。",
                 "一键容器化部署，多个角色知识库独立隔离。"])
    add_bullets(s, 7.0, 1.35, 5.3, 2.5, "后续可以继续优化的方向",
                ["跨页表格自动拼接：当一张表格跨越 PDF 多页时，自动合并成完整表格。",
                 "图片处理速度：图片 AI 分析比较慢，后续可以改为异步后台处理。",
                 "建立更全面的测试集：用更多真实问题评估系统效果，持续改进。",
                 "知识图谱轻量化更新：新上传文件后，只增量更新图谱，不用全部重建。"])
    add_textbox(s, 0.9, 5.65, 11.6, 0.55,
                "一句话总结：本项目让大模型具备「读懂私有文件、记住多轮对话、查找实体关系、基于证据回答」的完整能力。",
                18, TITLE, True, PP_ALIGN.CENTER)
    add_footer(s, 13)
    add_transition(s, "push_left", 500)

    # ------------------------------------------------------------------ #
    # SLIDE 14  结束页                                                      #
    # ------------------------------------------------------------------ #
    s = blank_slide(prs)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
    add_textbox(s, 0.8, 2.1,  11.8, 0.75, "谢谢聆听", 42, TITLE, True, PP_ALIGN.CENTER)
    add_textbox(s, 0.8, 3.05, 11.8, 0.45, "欢迎各位老师提问与批评指正", 20, ACCENT, True, PP_ALIGN.CENTER)
    add_textbox(s, 0.8, 4.15, 11.8, 0.35, "基于 RAG 的角色扮演问答系统  |  2309A_nlp-李学远", 14, MUTED, False, PP_ALIGN.CENTER)
    add_footer(s, 14)
    add_transition(s, "fade", 800)

    prs.save(PPT)
    print(f"saved:  {PPT}")
    print(f"backup: {BACKUP}")
    print(f"slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_ppt()
